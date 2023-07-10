import cgi
import os
import PyPDF2
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

opq_file = open('LS_OPQ.pdf', 'rb')
verify_file = open('LS_Verify.pdf', 'rb')

#this is input from the checked boxes on website
def extract_dimensions(form):
    dimensions = []
    if 'dimension' in form:
        dimensions = form.getlist('dimension')
    return dimensions



def handle_form_submission(form):
    # Save the uploaded files to the server
    opq_file = form['opq_file']
    verify_file = form['verify_file']
    opq_file_path = os.path.join('./uploads', opq_file.filename)
    verify_file_path = os.path.join('./uploads', verify_file.filename)
    opq_file.save(opq_file_path)
    verify_file.save(verify_file_path)

    # Extract the dimensions from the form data
    dimensions = extract_dimensions(form)

    # Run your existing code with the file paths and dimensions
    run(opq_file_path, verify_file_path, dimensions)

    # Return a success message
    return 'Conversion successful'

def selecting_dimmensions():
	my_list = ['Besluta och starta aktiviteter', 'Leda och följa upp', 'Arbeta med människor',
	            'Skapa relationer och nätverk', 'Uppnå personliga arbetsmål']
	return my_list

def extract_all_info_opq():
	pdf_reader = PyPDF2.PdfReader(opq_file)
	text = ""
	for page in range(len(pdf_reader.pages)):
		page_obj = pdf_reader.pages[page]
		text += page_obj.extract_text()

	marker = "1. Leda och Ta beslut 12345"

	lines = text.split("\n")
	index = next((i for i, line in enumerate(lines) if marker in line), None)

	if index is not None:
		trimmed_string = "\n".join(lines[index:])

	# print(trimmed_string)
	index = trimmed_string.find(
		'1.1 Besluta  och starta aktiviteterTar ansvar för åtgärder, projekt och  personal; tar initiativ  och arbetar')

	if index != -1:
		return trimmed_string[:index]
	else:
		return trimmed_string

def extract_all_verify():
	pdf_reader = PyPDF2.PdfReader(verify_file)
	text = ""
	for page in range(len(pdf_reader.pages)):
		page_obj = pdf_reader.pages[page]
		text += page_obj.extract_text()

	return text

def string_to_list_opq():
	list_of_words = ['DeDN', 'Spacer', '1', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '2.2', '3.1', '3.2', '3.3', '4.1',
	                 '4.2', '4.3', '5.1', '5.2', '5.3', '6.1', '6.2', '6.3', '7.1', '7.2', '8.1', '8.2', '\uf0fc\uf0fc',
	                 '\uf0fc', '\uf06c\uf06c', '\uf06c', '\uf0fb\uf0fb', '\uf0fb', '©']
	trimmed_string = extract_all_info_opq()
	len_list = trimmed_string.split("\n")
	list = trimmed_string.split("\n")
	remove_list = []
	final_list = []

	for i in range(len(len_list)):

		if not any(list[i].startswith(word) for word in list_of_words):
			joined_string = list[i - 1] + list[i]
			list[i - 1] = joined_string
			remove_list.append(list[i])

	for word in list:
		if word not in remove_list:
			final_list.append(word)

	modified_list = []

	for string in final_list:
		while "  " in string:
			string = string.replace("  ", " ")
		modified_list.append(string)

	return modified_list

def string_to_list_verify():
	long_string = extract_all_verify()
	long_list = long_string.split("\n")

	name = long_list[1]
	generell = long_list[12].split('G+')[1]
	induktiv = long_list[13].split('Induktivt resonemang')[1]
	numerisk = long_list[14].split('Numeriska färdigheter')[1]
	deduktiv = long_list[15].split('Deduktivt resonemang')[1]

	induktiv_text_list = long_list[38:41]
	numerisk_text_list = long_list[53:56]
	deduktiv_text_list = long_list[64:70]

	induktiv_string = ''.join(induktiv_text_list)
	numerisk_string = ''.join(numerisk_text_list)
	deduktiv_string = ''.join(deduktiv_text_list)

	all_percent_list = [generell, induktiv, numerisk, deduktiv]
	all_percent_list = [s.replace(' ', '') for s in all_percent_list]
	all_text_list = [induktiv_string, numerisk_string, deduktiv_string]

	modified_list = []

	for string in all_text_list:
		while "  " in string:
			string = string.replace("  ", " ")
		modified_list.append(string)

	return all_percent_list, modified_list, name

def get_dimmension(index, number):
	chosen_list = []
	full_list = string_to_list_opq()
	dimmension = full_list[index].split(" ", 1)
	chosen_list.append(dimmension[1])

	for n in range(index + 1, index + number + 1):
		a = full_list[n].split(" ", 1)
		chosen_list.append(a)

	return chosen_list

def get_chosen_dimmension_list():
	list_of_dimmenssions = [['Besluta och starta aktiviteter', 1, 4], ['Leda och följa upp', 6, 4],
	                        ['Arbeta med människor', 12, 5], ['Stå fast vid principer och värderingar', 18, 2],
	                        ['Skapa relationer och nätverk', 22, 4], ['Övertala och påverka', 27, 5],
	                        ['Presentera och kommunicera information', 33, 4], ['Skriva och rapportera', 41, 4],
	                        ['Tillämpa expertis och teknologi', 46, 3], ['Analysera', 50, 3], ['Lära och utforska', 55, 4],
	                        ['Skapa och uppfinna', 60, 4], ['Formulera strategier och koncept', 65, 4],
	                        ['Planera och organisera', 72, 4], ['Leverera resultat och uppfylla kundförväntningar', 77, 4],
	                        ['Följa instruktioner och procedurer', 82, 3], ['Anpassa och reagera på förändring', 87, 4],
	                        ['Hantera krav och motgångar', 92, 4], ['Uppnå personliga arbetsmål', 99, 4],
	                        ['Företagaranda och kommersiellt  tänkande', 104, 3]]

	string_list = []
	index_list = []
	for i in range(len(list_of_dimmenssions)):
		d = get_dimmension(list_of_dimmenssions[i][1], list_of_dimmenssions[i][2])
		if list_of_dimmenssions[i][0] in selecting_dimmensions():
			index_list.append(list_of_dimmenssions[i][2])
			string_list.append(d)


	return string_list, index_list

def make_slide_8_to_12():
	presentation = Presentation("Fördjupad bedömning - mall.pptx")
	slide_index = 7
	dimmension_list, index_list = get_chosen_dimmension_list()
	only_dimension_list = selecting_dimmensions()

	for i in range(len(only_dimension_list)):

		# Ensure the slide index is valid
		if slide_index < len(presentation.slides):
			slide = presentation.slides[slide_index]

			textbox_index = 0  # Modify the first textbox
			if len(slide.shapes) > textbox_index:
				shape = slide.shapes[textbox_index]
				if shape.has_text_frame:
					text_frame = shape.text_frame

					# Modify the text content
					text_frame.text = only_dimension_list[i]

					# Modify the font size
					for paragraph in text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(24)
							run.font.bold = False

			textbox_index = 1  # Modify the second textbox
			if len(slide.shapes) > textbox_index:
				shape = slide.shapes[textbox_index]
				if shape.has_text_frame:
					text_frame = shape.text_frame

					# Modify the text content
					text_frame.text = only_dimension_list[i]

					# Modify the font size
					for paragraph in text_frame.paragraphs:
						for run in paragraph.runs:
							run.font.size = Pt(16)
							run.font.bold = False

							# Modify the font color
							font_color = RGBColor(0, 0, 0)  # Customize the RGB color values as needed
							run.font.color.rgb = font_color

			# Iterate through the shapes to find the table
			for shape in slide.shapes:
				if shape.has_table:
					table = shape.table

					# Check if the table has at least four rows and two columns
					if len(table.rows) >= 5 and len(table.columns) > 1:
						# Modify the first and second columns
						texts_first_column = []
						texts_second_column =[]
						for b in range(index_list[i]):
							texts_first_column.append(dimmension_list[i][b+1][0])
							texts_second_column.append(dimmension_list[i][b+1][1])

						for n in range(index_list[i]):  # Modify the current row and the next three rows
							#current_row_index = n  # Skip the header row

							# Modify cell in the first column
							cell_first_column = table.cell(n, 0)
							if cell_first_column.text:
								text_frame_first_column = cell_first_column.text_frame
								paragraph_first_column = text_frame_first_column.paragraphs[0]
								run_first_column = paragraph_first_column.runs[0]

								run_first_column.text = texts_first_column[n]
								run_first_column.font.bold = False
								run_first_column.font.size = Pt(14)

								if text_frame_first_column.text.startswith('\uf0fc'):
									run_first_column.font.color.rgb = RGBColor(0, 255, 0)
								elif text_frame_first_column.text.startswith('\uf0fb'):
									run_first_column.font.color.rgb = RGBColor(255, 0, 0)
								else:
									run_first_column.font.color.rgb = RGBColor(0, 0, 0)

							# Modify cell in the second column
							cell_second_column = table.cell(n, 1)
							if cell_second_column.text:
								text_frame_second_column = cell_second_column.text_frame
								paragraph_second_column = text_frame_second_column.paragraphs[0]
								run_second_column = paragraph_second_column.runs[0]

								run_second_column.text = texts_second_column[n]
								run_second_column.font.bold = False
								run_second_column.font.size = Pt(14)
								run_second_column.font.color.rgb = RGBColor(0, 0, 0)
		slide_index = 8 + i


	modified_presentation_path = "modified_presentation.pptx"
	presentation.save(modified_presentation_path)

def make_slide_4_and_7_OPQ(slide_index, font_size):
	presentation = Presentation("modified_presentation.pptx")


	if slide_index < len(presentation.slides):
		slide = presentation.slides[slide_index]

		# Method: Find the table based on cell values
		target_table_cells = [
			["Test"],
			["Test"],
			["Test"],
			["Test"],
			["Test"],
		]
		new_texts = selecting_dimmensions()

		for shape in slide.shapes:
			if shape.has_table:
				table = shape.table

				# Check if the table has the desired number of rows and columns
				if len(table.rows) >= len(target_table_cells) and len(table.columns) >= 1:
					found_table = True

					# Check if the cell values match the target_table_cells
					for row_index, row_cells in enumerate(target_table_cells):
						for col_index, cell_value in enumerate(row_cells):
							cell = table.cell(row_index, col_index)
							if cell.text != cell_value:
								found_table = False
								break

						if not found_table:
							break

					if found_table:
						# Modify the content of the matched table
						for row_index, row_cells in enumerate(target_table_cells):
							for col_index, cell_value in enumerate(row_cells):
								cell = table.cell(row_index, col_index)
								text_frame = cell.text_frame

								if text_frame.paragraphs:
									# Access the first paragraph in the text frame
									paragraph = text_frame.paragraphs[0]

									# Store the existing font properties
									font_name = paragraph.runs[0].font.name


									# Clear existing content if needed
									paragraph.clear()

									# Create a new run and set the properties
									run = paragraph.add_run()
									run.text = new_texts[row_index]  # Replace with the new text from the list
									run.font.name = font_name
									run.font.size = Pt(font_size)
									run.font.color.rgb = RGBColor(0, 0, 0)
									run.font.bold = False

						break

	# Save the modified presentation to a new file
	modified_presentation_path = "modified_presentation.pptx"
	presentation.save(modified_presentation_path)

def make_slide_4_verify():
	presentation = Presentation("modified_presentation.pptx")
	slide_index = 3
	percent_list, text_list, name = string_to_list_verify()

	if slide_index < len(presentation.slides):
		slide = presentation.slides[slide_index]

		# Method: Find the table based on cell values
		target_table_cells = [
			["G"],
			["I"],
			["N"],
			["D"],
		]

		new_texts = percent_list
		color = RGBColor(0, 0, 0)  # Specify your own RGB color values

		for shape in slide.shapes:
			if shape.has_table:
				table = shape.table

				# Check if the table has the desired number of rows and columns
				if len(table.rows) >= len(target_table_cells) and len(table.columns) >= 2:
					found_table = True

					# Check if the cell values match the target_table_cells
					for row_index, row_cells in enumerate(target_table_cells):
						for col_index, cell_value in enumerate(row_cells):
							cell = table.cell(row_index, col_index + 1)
							if cell.text != cell_value:
								found_table = False
								break

						if not found_table:
							break

					if found_table:
						# Modify the content of the matched table
						for row_index, row_cells in enumerate(target_table_cells):
							col_index = 1  # Change the column index to 1 (second column)
							cell = table.cell(row_index, col_index)

							if cell.text_frame.paragraphs:
								# Access the existing text frame and paragraph in the cell
								text_frame = cell.text_frame
								paragraph = text_frame.paragraphs[0]

								# Clear existing content if needed
								paragraph.clear()

								# Create a new run and set the properties
								run = paragraph.add_run()
								run.text = new_texts[row_index]  # Replace with the new text from the list

								# Copy the formatting from the first column's text run
								source_cell = table.cell(row_index, 0)
								if source_cell.text_frame.paragraphs:
									source_paragraph = source_cell.text_frame.paragraphs[0]
									source_run = source_paragraph.runs[0]

									# Set the font properties of the second column based on the first column
									run.font.name = source_run.font.name
									run.font.size = source_run.font.size
									run.font.bold = source_run.font.bold
									run.font.italic = source_run.font.italic
									run.font.underline = source_run.font.underline

									# Set the color of the second column
									run.font.color.rgb = color

						break



	modified_presentation_path = "modified_presentation.pptx"
	presentation.save(modified_presentation_path)

def make_slide_16():
	presentation = Presentation("modified_presentation.pptx")
	slide_index = 15
	percent_list, text_list, name = string_to_list_verify()

	if slide_index < len(presentation.slides):
		slide = presentation.slides[slide_index]

		# Method: Find the table based on cell values
		target_table_cells = [
			["G"],
			["I"],
			["N"],
			["D"],
		]

		new_texts = percent_list

		for shape in slide.shapes:
			if shape.has_table:
				table = shape.table

				# Check if the table has the desired number of rows and columns
				if len(table.rows) >= len(target_table_cells) and len(table.columns) >= 2:
					found_table = True

					# Check if the cell values match the target_table_cells
					for row_index, row_cells in enumerate(target_table_cells):
						for col_index, cell_value in enumerate(row_cells):
							cell = table.cell(row_index, col_index + 1)
							if cell.text != cell_value:
								found_table = False
								break

						if not found_table:
							break

					if found_table:
						# Modify the content of the matched table
						for row_index, row_cells in enumerate(target_table_cells):
							col_index = 1  # Change the column index to 1 (second column)
							cell = table.cell(row_index, col_index)
							text_frame = cell.text_frame

							if text_frame.paragraphs:
								# Access the first paragraph in the text frame
								paragraph = text_frame.paragraphs[0]

								# Store the existing font properties
								font_name = paragraph.runs[0].font.name

								# Clear existing content if needed
								paragraph.clear()

								# Create a new run and set the properties
								run = paragraph.add_run()
								run.text = new_texts[row_index]  # Replace with the new text from the list
								run.font.name = font_name
								run.font.size = Pt(20)
								run.font.color.rgb = RGBColor(0, 0, 0)
								run.font.bold = False

								shape.fill.solid()  # Set the fill type to solid color
								shape.fill.fore_color.rgb = RGBColor(255, 0, 0)

						break

	modified_presentation_path = "modified_presentation.pptx"
	presentation.save(modified_presentation_path)

def make_slide_17(inputs, ability):
	presentation = Presentation("modified_presentation.pptx")
	percent_list, text_list, name = string_to_list_verify()
	slide_index = 16
	i = inputs

	if slide_index < len(presentation.slides):
		slide = presentation.slides[slide_index]

		# Method: Find the table based on cell values
		target_table_cells = [
			[ability],
			["Test"],
		]
		new_texts = text_list[i]
		percents = percent_list[i+1]

		for shape in slide.shapes:
			if shape.has_table:
				table = shape.table

				# Check if the table has the desired number of rows and columns
				if len(table.rows) >= len(target_table_cells) and len(table.columns) >= 1:
					found_table = True

					# Check if the cell values match the target_table_cells
					for row_index, row_cells in enumerate(target_table_cells):
						for col_index, cell_value in enumerate(row_cells):
							cell = table.cell(row_index, col_index)
							if cell.text != cell_value:
								found_table = False
								break

						if not found_table:
							break

					if found_table:
						# Modify the content of the matched table

						cell = table.cell(1, 0)
						text_frame = cell.text_frame

						if text_frame.paragraphs:
							# Access the first paragraph in the text frame
							paragraph = text_frame.paragraphs[0]

							# Store the existing font properties
							font_name = paragraph.runs[0].font.name

							# Clear existing content if needed
							paragraph.clear()

							# Create a new run and set the properties
							run = paragraph.add_run()
							run.text = new_texts  # Replace with the new text from the list
							run.font.name = font_name
							run.font.size = Pt(10)
							run.font.color.rgb = RGBColor(0, 0, 0)
							run.font.bold = False

						cell = table.cell(0, 1)
						text_frame = cell.text_frame

						if text_frame.paragraphs:
							# Access the first paragraph in the text frame
							paragraph = text_frame.paragraphs[0]

							# Store the existing font properties
							font_name = paragraph.runs[0].font.name

							# Clear existing content if needed
							paragraph.clear()

							# Create a new run and set the properties
							run = paragraph.add_run()
							run.text = percents  # Replace with the new text from the list
							run.font.name = font_name
							run.font.size = Pt(17)
							run.font.color.rgb = RGBColor(0, 0, 0)


						break

	# Save the modified presentation to a new file
	modified_presentation_path = "modified_presentation.pptx"
	presentation.save(modified_presentation_path)

def run():
	make_slide_8_to_12()
	make_slide_4_and_7_OPQ(6, 11)
	make_slide_4_and_7_OPQ(3, 8)
	make_slide_17(0, 'Induktiv förmåga')
	make_slide_17(1,'Numerisk förmåga')
	make_slide_17(2,'Deduktiv förmåga')
	make_slide_16()
	make_slide_4_verify()



def main():
    form = cgi.FieldStorage()

    # Check if the form has been submitted
    if 'opq_file' in form and 'verify_file' in form:
        message = handle_form_submission(form)
    else:
        message = ''

    # Print the message as a response
    print('Content-type: text/html\n')
    print(message)

# Run the main function when the script is executed
if __name__ == '__main__':
    main()
